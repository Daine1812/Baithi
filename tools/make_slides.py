#!/usr/bin/env python3
from __future__ import annotations

import argparse
import sys
import os
from typing import List, Tuple, Optional

try:
    import cv2  # type: ignore
except Exception:
    cv2 = None  # Optional dependency; script will work without it

from PIL import Image
import numpy as np
import pytesseract
from pytesseract import TesseractError
try:
    import easyocr  # type: ignore
except Exception:
    easyocr = None

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor


def parse_args(argv: Optional[List[str]] = None) -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Generate a PPTX from one or more images by OCR-ing Vietnamese text.",
    )
    parser.add_argument(
        "images",
        nargs="+",
        help="Paths to input image files (e.g., yeucau.jpg, de bai.jpg)",
    )
    parser.add_argument(
        "--output",
        "-o",
        default="slides_from_images.pptx",
        help="Output PPTX file path",
    )
    parser.add_argument(
        "--lang",
        "-l",
        default="vie",
        help="Tesseract language code to use (default: vie)",
    )
    parser.add_argument(
        "--fallback-lang",
        default="eng",
        help="Fallback language if the primary language is unavailable (default: eng)",
    )
    parser.add_argument(
        "--no-preprocess",
        action="store_true",
        help="Disable image preprocessing before OCR",
    )
    parser.add_argument(
        "--title-from",
        choices=["first-line", "filename"],
        default="first-line",
        help="How to derive the slide title (default: first-line)",
    )
    parser.add_argument(
        "--font-name",
        default="DejaVu Sans",
        help="Font to use for text (default: DejaVu Sans)",
    )
    parser.add_argument(
        "--title-size",
        type=int,
        default=40,
        help="Font size for titles in points (default: 40)",
    )
    parser.add_argument(
        "--bullet-size",
        type=int,
        default=24,
        help="Font size for bullet points in points (default: 24)",
    )
    parser.add_argument(
        "--accent-color",
        default="#1f77b4",
        help="Hex color for titles and accents (default: #1f77b4)",
    )
    parser.add_argument(
        "--wide",
        action="store_true",
        help="Use 16:9 (widescreen) slide size",
    )
    return parser.parse_args(argv)


def hex_to_rgb_color(hex_str: str) -> RGBColor:
    hex_str = hex_str.strip().lstrip("#")
    if len(hex_str) != 6:
        raise ValueError("accent-color must be a 6-digit hex string, e.g., #1f77b4")
    r = int(hex_str[0:2], 16)
    g = int(hex_str[2:4], 16)
    b = int(hex_str[4:6], 16)
    return RGBColor(r, g, b)


def load_image(image_path: str) -> Image.Image:
    return Image.open(image_path)


def preprocess_image_for_ocr(pil_image: Image.Image) -> Image.Image:
    if cv2 is None:
        gray = pil_image.convert("L")
        # Simple binarization fallback without OpenCV
        return gray.point(lambda x: 0 if x < 180 else 255)

    np_img = np.array(pil_image)
    if len(np_img.shape) == 3:
        gray = cv2.cvtColor(np_img, cv2.COLOR_BGR2GRAY)
    else:
        gray = np_img
    # Adaptive threshold to improve OCR on varied backgrounds
    thresh = cv2.adaptiveThreshold(
        gray,
        255,
        cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
        cv2.THRESH_BINARY,
        35,
        11,
    )
    return Image.fromarray(thresh)


def run_ocr(pil_image: Image.Image, primary_lang: str, fallback_lang: str) -> str:
    # Try Tesseract first
    try:
        return pytesseract.image_to_string(pil_image, lang=primary_lang)
    except (TesseractError, pytesseract.TesseractNotFoundError):
        pass
    # Fallback to EasyOCR if available
    if easyocr is not None:
        try:
            def map_lang(l: str) -> str:
                l = (l or "").lower()
                if l in ("vie", "vi_vn", "vi-vn", "vn", "vietnamese"):
                    return "vi"
                if l in ("eng", "en_us", "en-us", "english"):
                    return "en"
                # best-effort: use first two letters
                return l[:2] if len(l) >= 2 else l

            langs = []
            if primary_lang:
                langs.append(map_lang(primary_lang))
            if fallback_lang:
                fb = map_lang(fallback_lang)
                if fb not in langs:
                    langs.append(fb)
            reader = easyocr.Reader(langs, gpu=False)
            result = reader.readtext(np.array(pil_image), detail=0)
            return "\n".join(result)
        except Exception:
            pass
    # Last resort: try Tesseract with fallback lang, else return empty
    try:
        return pytesseract.image_to_string(pil_image, lang=fallback_lang)
    except Exception:
        return ""


def normalize_line(text_line: str) -> str:
    stripped = text_line.strip()
    # Remove common bullet characters and numbering
    bullets = ["•", "-", "–", "—", "*", "·"]
    for b in bullets:
        if stripped.startswith(b):
            stripped = stripped[len(b):].strip()
            break
    # Remove simple numeric prefixes like "1.", "1)"
    if len(stripped) >= 2 and stripped[0].isdigit() and stripped[1] in ".)":
        stripped = stripped[2:].strip()
    return stripped


def parse_text_to_title_and_bullets(
    ocr_text: str, title_from: str, filename: str
) -> Tuple[str, List[str]]:
    lines = [ln for ln in (ln.strip("\r").rstrip() for ln in ocr_text.splitlines())]
    non_empty = [ln for ln in lines if ln.strip()]
    if title_from == "filename" or not non_empty:
        title = os.path.splitext(os.path.basename(filename))[0]
    else:
        title = normalize_line(non_empty[0])

    bullet_candidates = non_empty[1:] if title_from == "first-line" and len(non_empty) > 1 else non_empty
    bullets: List[str] = []
    for raw in bullet_candidates:
        normalized = normalize_line(raw)
        if not normalized:
            continue
        if normalized.lower() == title.lower():
            continue
        bullets.append(normalized)

    # Deduplicate while preserving order
    seen = set()
    unique_bullets: List[str] = []
    for item in bullets:
        if item not in seen:
            unique_bullets.append(item)
            seen.add(item)

    return title, unique_bullets


def set_text_run_font(run, font_name: str, size_pt: int, color: RGBColor | None = None):
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    if color is not None:
        run.font.color.rgb = color


def add_title_and_bullets_slide(
    prs: Presentation,
    title: str,
    bullets: List[str],
    font_name: str,
    title_size_pt: int,
    bullet_size_pt: int,
    accent_color: RGBColor,
):
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)

    title_shape = slide.shapes.title
    title_shape.text = title
    for paragraph in title_shape.text_frame.paragraphs:
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
        if paragraph.runs:
            for run in paragraph.runs:
                set_text_run_font(run, font_name, title_size_pt, accent_color)

    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    if not bullets:
        p = body.paragraphs[0]
        run = p.add_run()
        run.text = ""
        set_text_run_font(run, font_name, bullet_size_pt)
        return

    first = True
    for bullet in bullets:
        if first:
            p = body.paragraphs[0]
            first = False
        else:
            p = body.add_paragraph()
        p.level = 0
        run = p.add_run()
        run.text = bullet
        set_text_run_font(run, font_name, bullet_size_pt)


def build_presentation(
    images: List[str],
    output_path: str,
    lang: str,
    fallback_lang: str,
    use_preprocess: bool,
    title_from: str,
    font_name: str,
    title_size_pt: int,
    bullet_size_pt: int,
    accent_color_hex: str,
    wide: bool,
):
    prs = Presentation()
    if wide:
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

    accent_color = hex_to_rgb_color(accent_color_hex)

    for image_path in images:
        if not os.path.exists(image_path):
            print(f"[warn] File not found, skipping: {image_path}", file=sys.stderr)
            continue

        image = load_image(image_path)
        processed = (
            preprocess_image_for_ocr(image) if use_preprocess else image
        )
        text = run_ocr(processed, lang, fallback_lang)
        title, bullets = parse_text_to_title_and_bullets(text, title_from, image_path)
        add_title_and_bullets_slide(
            prs,
            title,
            bullets,
            font_name,
            title_size_pt,
            bullet_size_pt,
            accent_color,
        )

    prs.save(output_path)
    print(f"[ok] Saved presentation to: {output_path}")


def main(argv: Optional[List[str]] = None) -> int:
    args = parse_args(argv)
    build_presentation(
        images=args.images,
        output_path=args.output,
        lang=args.lang,
        fallback_lang=args.fallback_lang,
        use_preprocess=not args.no_preprocess,
        title_from=args.title_from,
        font_name=args.font_name,
        title_size_pt=args.title_size,
        bullet_size_pt=args.bullet_size,
        accent_color_hex=args.accent_color,
        wide=args.wide,
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
