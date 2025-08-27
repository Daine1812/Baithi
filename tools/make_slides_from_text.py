#!/usr/bin/env python3
from __future__ import annotations

import argparse
from typing import List, Tuple, Union

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor


Bullet = Union[str, Tuple[str, List[str]]]


def set_font(run, font_name: str, size_pt: int, color: RGBColor | None = None):
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    if color is not None:
        run.font.color.rgb = color


def add_title_and_bullets(
    prs: Presentation,
    title: str,
    bullets: List[Bullet],
    font_name: str,
    title_size: int,
    bullet_size: int,
    accent_color: RGBColor,
):
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)

    title_shape = slide.shapes.title
    title_shape.text = title
    for paragraph in title_shape.text_frame.paragraphs:
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
        for run in paragraph.runs:
            set_font(run, font_name, title_size, accent_color)

    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    first_para = True
    for item in bullets:
        if isinstance(item, tuple):
            parent, children = item
            p = body.paragraphs[0] if first_para else body.add_paragraph()
            first_para = False
            p.level = 0
            run = p.add_run()
            run.text = parent
            set_font(run, font_name, bullet_size)

            for child in children:
                c = body.add_paragraph()
                c.level = 1
                run = c.add_run()
                run.text = child
                set_font(run, font_name, bullet_size)
        else:
            p = body.paragraphs[0] if first_para else body.add_paragraph()
            first_para = False
            p.level = 0
            run = p.add_run()
            run.text = item
            set_font(run, font_name, bullet_size)


def hex_to_rgb_color(hex_str: str) -> RGBColor:
    hex_str = hex_str.strip().lstrip("#")
    r = int(hex_str[0:2], 16)
    g = int(hex_str[2:4], 16)
    b = int(hex_str[4:6], 16)
    return RGBColor(r, g, b)


def build_studymate_ai_presentation(output_path: str):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    font_name = "DejaVu Sans"
    title_size = 42
    bullet_size = 24
    accent = hex_to_rgb_color("#1f77b4")

    # Slide 1: Title
    add_title_and_bullets(
        prs,
        title="StudyMate AI – Trợ lý học tập thông minh",
        bullets=[
            "Ý tưởng khởi nghiệp CNTT",
        ],
        font_name=font_name,
        title_size=title_size,
        bullet_size=bullet_size,
        accent_color=accent,
    )

    # Slide 2: Ý tưởng khởi nghiệp
    add_title_and_bullets(
        prs,
        title="1. Ý tưởng khởi nghiệp",
        bullets=[
            "Tên dự án: StudyMate AI",
            "Sản phẩm: Ứng dụng mobile & web giúp sinh viên học tập thông minh bằng AI",
            ("Lý do chọn ý tưởng:", [
                "Khó quản lý thời gian học, làm bài tập, ghi nhớ",
                "Học online nhưng thiếu công cụ cá nhân hóa",
            ]),
            ("Điểm khác biệt:", [
                "AI gợi ý lộ trình học cá nhân",
                "Tóm tắt bài giảng, gợi ý flashcard, tạo quiz",
                "Chatbot giải thích kiến thức như gia sư ảo",
            ]),
        ],
        font_name=font_name,
        title_size=title_size,
        bullet_size=bullet_size,
        accent_color=accent,
    )

    # Slide 3: Khách hàng mục tiêu & vấn đề
    add_title_and_bullets(
        prs,
        title="2. Khách hàng mục tiêu & vấn đề",
        bullets=[
            ("Chân dung khách hàng:", [
                "Sinh viên đại học, cao đẳng",
                "Học sinh THPT chuẩn bị thi",
            ]),
            ("Nhu cầu/vấn đề:", [
                "Khó quản lý lịch học, bài tập",
                "Thiếu công cụ tóm tắt nhanh, học hiệu quả",
            ]),
            "Minh chứng: 80% sinh viên muốn app tóm tắt và ôn tập nhanh",
        ],
        font_name=font_name,
        title_size=title_size,
        bullet_size=bullet_size,
        accent_color=accent,
    )

    # Slide 4: Giải pháp & MVP
    add_title_and_bullets(
        prs,
        title="3. Giải pháp & MVP",
        bullets=[
            ("Giải pháp – StudyMate AI hỗ trợ:", [
                "Tải PDF/Word → AI tóm tắt bullet points",
                "Sinh flashcard & quiz ôn tập tự động",
                "Chatbot hỏi–đáp như trợ giảng ảo",
                "Quản lý lịch học, nhắc deadline",
            ]),
            "MVP: Bản web demo với Tóm tắt tài liệu & Tạo quiz",
            "Đo lường: số lượt tải, giờ học trung bình/ngày",
        ],
        font_name=font_name,
        title_size=title_size,
        bullet_size=bullet_size,
        accent_color=accent,
    )

    # Slide 5: Business Model Canvas – StudyMate AI
    add_title_and_bullets(
        prs,
        title="4. Business Model Canvas – StudyMate AI",
        bullets=[
            ("1) Đối tác chính (Key Partners):", [
                "OpenAI, HuggingFace, Google AI (AI/ML)",
                "Trường đại học, trung tâm giáo dục",
                "Đối tác thanh toán: Momo, ZaloPay, VNPay",
                "Startup EdTech, nhà xuất bản tài liệu",
            ]),
            ("2) Hoạt động chính (Key Activities):", [
                "Phát triển & duy trì app (mobile/web)",
                "Xây dựng/huấn luyện mô hình AI (tóm tắt, quiz, chatbot)",
                "Marketing online/offline tại trường học",
                "CSKH & hỗ trợ kỹ thuật",
            ]),
            ("3) Giá trị cốt lõi (Value Proposition):", [
                "Học thông minh hơn, tiết kiệm thời gian",
                "Cá nhân hóa lộ trình, ôn tập hiệu quả",
                "Trợ lý ảo AI: tóm tắt, quiz, flashcard",
                "Khác biệt: tự động hóa – cá nhân hóa – tương tác như gia sư",
            ]),
            ("4) Quan hệ khách hàng (Customer Relationships):", [
                "Miễn phí + nâng cấp Premium",
                "Hỗ trợ chatbot 24/7, cộng đồng Facebook/Zalo",
                "Gamification: tích điểm đổi thưởng",
                "Email/SMS nhắc lịch học, deadline",
            ]),
            ("5) Phân khúc khách hàng (Customer Segments):", [
                "Sinh viên đại học, cao đẳng",
                "Học sinh THPT chuẩn bị thi",
                "Người đi làm muốn học thêm",
            ]),
            ("6) Kênh phân phối (Channels):", [
                "App Store, Google Play",
                "Website chính thức",
                "MXH: Facebook, TikTok, YouTube",
                "Hợp tác CLB sinh viên, trung tâm gia sư",
            ]),
            ("7) Nguồn lực chính (Key Resources):", [
                "Đội ngũ dev & chuyên gia AI",
                "Hạ tầng cloud: AWS, GCP",
                "Dữ liệu học tập (giáo trình, đề thi)",
                "Vốn khởi nghiệp/đầu tư",
            ]),
            ("8) Cơ cấu chi phí (Cost Structure):", [
                "Phát triển ứng dụng & server cloud",
                "Nhân sự: dev, AI, marketing",
                "Marketing & quảng cáo",
                "Bản quyền AI/API",
            ]),
            ("9) Dòng doanh thu (Revenue Streams):", [
                "Gói Premium: 99k/tháng (AI nâng cao, flashcard không giới hạn)",
                "Quảng cáo (phiên bản free)",
                "B2B: Giải pháp AI cho trường học/trung tâm",
                "Khóa học mini tích hợp trong app",
            ]),
        ],
        font_name=font_name,
        title_size=title_size,
        bullet_size=bullet_size,
        accent_color=accent,
    )

    # Slide 6: Go-to-market & Marketing
    add_title_and_bullets(
        prs,
        title="5. Tiếp cận thị trường & marketing",
        bullets=[
            ("Kênh phân phối:", [
                "App Store, Google Play, website",
            ]),
            ("Marketing:", [
                "Hợp tác CLB sinh viên, phát demo miễn phí",
                "Quảng cáo Facebook, TikTok, YouTube",
                "Mini game: Ôn thi cùng AI",
            ]),
        ],
        font_name=font_name,
        title_size=title_size,
        bullet_size=bullet_size,
        accent_color=accent,
    )

    # Slide 7: Tài chính sơ bộ
    add_title_and_bullets(
        prs,
        title="6. Phân tích tài chính sơ bộ",
        bullets=[
            ("Chi phí ban đầu:", [
                "Phát triển ứng dụng: 50 triệu",
                "Marketing thử nghiệm: 20 triệu",
            ]),
            "Nguồn thu: Gói Premium & quảng cáo trong app",
            ("Dự kiến lợi nhuận:", [
                "1.000 người dùng trả phí → 99 triệu/tháng",
                "Hòa vốn sau 6 tháng",
            ]),
        ],
        font_name=font_name,
        title_size=title_size,
        bullet_size=bullet_size,
        accent_color=accent,
    )

    # Slide 8: Lộ trình 1 năm
    add_title_and_bullets(
        prs,
        title="7. Kế hoạch phát triển 1 năm",
        bullets=[
            "Q1: Ra mắt MVP, test với 100 sinh viên",
            "Q2: Ra mắt trên Google Play, đạt 10.000 user",
            "Q3: Nâng cấp AI Chatbot theo môn học",
            "Q4: Hợp tác trường học, mở rộng bản tiếng Anh",
        ],
        font_name=font_name,
        title_size=title_size,
        bullet_size=bullet_size,
        accent_color=accent,
    )

    prs.save(output_path)


def parse_args():
    parser = argparse.ArgumentParser(description="Sinh slide StudyMate AI từ nội dung có sẵn")
    parser.add_argument("--output", "-o", default="StudyMate_AI.pptx", help="Đường dẫn file PPTX đầu ra")
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    build_studymate_ai_presentation(args.output)
    print(f"[ok] Đã tạo: {args.output}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
